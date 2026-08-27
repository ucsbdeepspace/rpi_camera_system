"""One-off builder for docs/session_results_2026-08-18_pid_tuning.pptx.

Summarizes the day's closed-loop PID tuning work: the old (~210-235Hz
control-rate) tuning baseline, the Pi-side ROI change that raised
telemetry to ~465Hz and broke it, the ring-down resonance remeasurement,
the fresh Ki/Kp search, the D-term retry, and the notch-filter result.
Not meant to be a general reusable tool -- rerun by hand if the deck
needs regenerating.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
BULLET_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x2E, 0x5C, 0x8A)

RESULTS = "results/"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(1.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = BULLET_COLOR
    return slide


def _add_title_bar(slide, title):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.05), Inches(12.8), Inches(1.05))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1.5)


def _add_bullets(slide, bullets, left, top, width, height, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(size)
        p.font.color.rgb = BULLET_COLOR
        p.space_after = Pt(8)
    return box


def add_image_slide(prs, title, bullets, image_path, img_left=Inches(6.9)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, title)
    _add_bullets(slide, bullets, Inches(0.5), Inches(1.25), Inches(6.1), Inches(5.9))
    if image_path:
        slide.shapes.add_picture(image_path, img_left, Inches(1.25), width=Inches(6.0))
    return slide


def add_full_image_slide(prs, title, bullets, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, title)
    _add_bullets(slide, bullets, Inches(0.5), Inches(1.25), Inches(12.3), Inches(1.4))
    slide.shapes.add_picture(image_path, Inches(1.9), Inches(2.55), height=Inches(4.7))
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(
        prs,
        "FTA Closed-Loop PID Tuning",
        "Session summary — gain search at the old control rate, a Pi-side "
        "throughput change that broke it, resonance remeasurement, and the "
        "notch-filter attempt",
    )

    add_image_slide(
        prs,
        "Where we started: best clean P+I result (~210-235Hz control rate)",
        [
            "Kp=1.75 counts/px held fixed all session; Ki was the effective speed lever, not Kp.",
            "Ki=200: 141ms rise, 1.1% overshoot, 141ms settling — a single clean transition, the best confirmed result.",
            "Ki=400 went faster (63ms) but with real 15% overshoot — a genuine speed/overshoot tradeoff, not free.",
            "Raising Kp alone never helped: same or worse settling, more overshoot, at every Ki tried.",
        ],
        RESULTS + "fta_closed_loop_step_response_tuning_panel.png",
        img_left=Inches(6.7),
    )

    add_image_slide(
        prs,
        "A real Kp instability boundary was found (~6-6.5 counts/px)",
        [
            "Kp=5.5 stable, Kp=6.0 marginal (decaying ringing), Kp=6.5+ genuinely unstable — growing oscillation, never settles.",
            "Far below the naive “matched-to-plant-gain” estimate (~11 counts/px) — phase lag erodes margin, static gain alone isn't a safe guide.",
            "Pushing toward the boundary didn't buy speed either — the marginal Kp=6.0 case settled slower than the safely-stable Kp=5.5.",
            "All excursions stayed physically bounded by the firmware's DAC clamp regardless of how bad the tuning was.",
        ],
        RESULTS + "fta_closed_loop_step_response_stability_panel.png",
        img_left=Inches(6.5),
    )

    add_image_slide(
        prs,
        "Pi-side ROI change raised telemetry ~210Hz → ~465Hz — the old gains broke",
        [
            "Real throughput win on the Pi side (ROI + other changes) — telemetry now streams at ~440-475Hz, up from ~210-235Hz.",
            "Updated the PID's fixed sample-time constant to match (1/210s → 1/457.5s) and rebuilt.",
            "Kp=1.75 / Ki=200 — the best result from the previous slide — is now genuinely UNSTABLE: growing oscillation, not settling.",
            "Also jumped wildly in X when stepping Y — checked directly: the OTHER axis (cy) stays flat through the step, so this is the controlled axis itself misbehaving, not cross-axis coupling.",
        ],
        RESULTS + "fta_closed_loop_step_response_vcp_20260818T191117Z.png",
    )

    add_image_slide(
        prs,
        "Ruled out sample rate itself as the cause",
        [
            "Hypothesis: more updates/second gives the integral term more chances to accumulate before the plant's ~41ms delay reflects earlier corrections back.",
            "Direct test: throttled the control loop back to ~200Hz while telemetry stayed at ~465Hz (decoupling the two).",
            "Same Kp=1.75/Ki=200 — still badly unstable (1180% overshoot) even at the OLD control rate.",
            "This falsifies “rate caused it” cleanly — something else changed. Static DC gain was also checked and ruled out (0.094-0.095 vs ~0.09 px/count, unchanged).",
        ],
        RESULTS + "fta_closed_loop_step_response_465hz_telemetry_200hz_control_kp1750_ki200.png",
    )

    add_image_slide(
        prs,
        "Ring-down resonance test (amp-off / pulse / amp-off, free decay)",
        [
            "User's own idea: preload a setpoint with the amp OFF (no physical effect), pulse the amp briefly, then cut it and watch the actuator ring down under pure mechanical dynamics — independent of the control loop or camera-rate phase fitting.",
            "First two attempts used host arrival timestamps — turned out to be bucketed into ~15-16ms bursts by Windows thread scheduling, giving wrong answers (22Hz, then an 8.6Hz curve_fit that looked plausible but wasn't).",
            "Fixed by timestamping from the firmware's own tick counter instead of the host clock, and abandoning curve_fit for a model-free peak/trough-spacing measurement.",
            "Real result: 16 peaks + 17 troughs, mean spacing 26.0ms → 38.5Hz — well outside the 10-20Hz target band, and different from the original (now known-unreliable) 15.3Hz reading.",
        ],
        RESULTS + "fta_ringdown_peak_spacing_analysis.png",
        img_left=Inches(6.5),
    )

    add_image_slide(
        prs,
        "Fresh Ki search — Ki=19 looked clean, but under a leftover throttle",
        [
            "Bisected Ki with Kp fixed at 1.75: 15 → 2265ms, 18 → 1953ms, 19 → 1797ms — clean and monotonic, no instability anywhere in this range.",
            "Kp=2.5 and Kp=3.5 both unstable, even just holding a fixed setpoint — the real Kp boundary sits below 2.5.",
            "Caveat found LATER (see two slides ahead): this whole search — Ki=15/18/19, the Kp check, and the D-term retry that follows — was unknowingly run with the control loop still capped at ~200Hz by a diagnostic throttle from earlier rate-falsification testing that never got reverted.",
            "So “Ki=19 is the new best result” did not hold up once that throttle was found and removed — treat this slide as the story of how the search unfolded, not the final answer.",
        ],
        RESULTS + "fta_closed_loop_465hz_ki_search.png",
        img_left=Inches(6.7),
    )

    add_image_slide(
        prs,
        "D-term retried with a resonance-informed cutoff — still doesn't help",
        [
            "Retried on top of the (still-throttled) Ki=19 baseline with a 10Hz derivative filter cutoff, informed by the corrected 38.5Hz resonance.",
            "Kd=0.005 → 3434% overshoot, wrong direction entirely. Kd=0.001 (smallest representable) → still 1392% overshoot.",
            "Decisive result regardless of the throttle issue: tested across cutoffs 1-20Hz and Kd 0.001-0.05, combined with both Ki=19 and the old Ki=200 — every single combination made things worse than P+I alone.",
            "Root cause: PIDController.hpp's derivative filter is a single-pole low-pass — any cutoff low enough to reject the resonance also removes the useful signal; any cutoff high enough to keep signal also passes the resonance straight through.",
        ],
        RESULTS + "fta_closed_loop_step_response_465hz_kp1750_ki19_kd5_fc10.png",
    )

    add_image_slide(
        prs,
        "Caught: the whole afternoon's tuning ran under a leftover ~200Hz throttle",
        [
            "The earlier rate-falsification test (checking whether sample rate itself caused instability) added a temporary throttle that gated the control loop to ~200Hz while telemetry stayed at the real ~465Hz — and it was never removed.",
            "Every result since then — Ki=15/18/19, the Kp=2.5/3.5 check, both D-term retests, and the first notch-filter test — was run at that artificial ~200Hz, not the true rate this project has actually been streaming at.",
            "Fixed: removed the throttle, restored the PID's internal time-step to match the true ~457.5Hz rate, rebuilt, reflashed.",
            "Retested the “best result,” Kp=1.75/Ki=19, at the TRUE rate: badly unstable — chaotic oscillation even while just holding the baseline, before the step is even applied.",
        ],
        RESULTS + "fta_closed_loop_step_response_fullrate_kp1750_ki19_notch385_q3.png",
    )

    add_image_slide(
        prs,
        "Real Ki search at the true rate, notch active — Ki=15 is honestly the best point",
        [
            "Ki=15: clean, bounded convergence, no growth (this plot). Ki=20: bounded but a visible slow “beating” envelope. Ki=30: same beating, now visibly growing.",
            "Also retried Kp=2.5/Ki=15 with the notch on, hoping it would buy Kp headroom — still badly unstable, same chaotic signature.",
            "The notch does not move the Ki or Kp stability boundary in any measurable way — the clean/marginal/unstable transition lands in essentially the same place with or without it.",
            "Honest best working point at the true rate, notch on or off: Kp=1.75, Ki=15 — well short of the pre-ROI-change 141ms/1.1% result.",
        ],
        RESULTS + "fta_closed_loop_step_response_fullrate_kp1750_ki15_notch385_q3.png",
    )

    add_image_slide(
        prs,
        "Sine tracking at Ki=15 reveals a much bigger problem than slow settling",
        [
            "Step response only tells you settling time for one jump — the real mission needs continuous rejection of a moving disturbance, so tested sine tracking at the “safe” Ki=15.",
            "Result: tracking gain (T) only 10-16% across 1-10Hz — the actuator barely follows a target moving that fast.",
            "For a feedback loop, S = 1 − T (sensitivity = disturbance rejection). Low T at a frequency IS poor rejection at that frequency — not a separate concern.",
            "T≈0.10-0.16 means S≈0.84-0.90: a real 5-10Hz beacon wobble would still show up as position error, 84-90% uncorrected. “Safe” Ki=15 is stable but not remotely fast enough for the actual mission.",
        ],
        RESULTS + "fta_closed_loop_onboard_sine_ki15_10Hz_15umpp_v2.png",
    )

    add_image_slide(
        prs,
        "A properly-built control-rate throttle unexpectedly recovers the OLD fast behavior",
        [
            "Tested directly: does throttling the control loop back to ~200Hz (this time atomically pairing the throttle with the PID's assumed sample time, so they can't drift apart) recover the pre-ROI-change Kp=1.75/Ki=200 behavior?",
            "Result: yes — 163-168ms rise, ~2% overshoot, ~193-201ms settling, confirmed reproducible on an immediate repeat. Essentially matches the historical 141-297ms range.",
            "Contradicts an EARLIER same-day attempt at nominally the same throttled config, which was catastrophically unstable (1180% overshoot) — could not identify a concrete mechanism for the discrepancy. Flagged honestly as an open mystery, not a solved one.",
        ],
        RESULTS + "fta_closed_loop_step_response_throttle200_kp1750_ki200_repeat.png",
    )

    add_image_slide(
        prs,
        "Sine-validated the recovered config — real improvement, still not enough",
        [
            "Kp=1.75/Ki=200, throttled to 200Hz: tracking gain 20-34% across 5-20Hz (vs. 10-16% for Ki=15) — a real ~2-3x improvement, clean and stable at every tested frequency.",
            "But S is still 66-80% across the band — most of a real disturbance would still get through. Fast step settling and good disturbance rejection at 10-20Hz are NOT the same property.",
            "Consistent with the much earlier open-loop finding that the actuator/plant itself has real rolloff in this band — not something PID tuning alone fixes.",
        ],
        RESULTS + "fta_closed_loop_onboard_sine_throttle200_ki200_10Hz_15umpp.png",
    )

    add_image_slide(
        prs,
        "Real closed-loop delay measured directly: ~11.5ms, not the assumed ~41ms",
        [
            "Built a firmware primitive (pulse_step) that timestamps both the DAC step and the telemetry response on the SAME firmware clock — no host-clock timing anywhere in the measurement, avoiding every host-timing bug hit earlier this session.",
            "Result: 11.5ms mean delay, tight (std 1.7ms), 6/6 consistent trials — much smaller than the ~41ms figure this project had been reasoning from (that number was from an older, slower VCP-relay test method, not the real I2C-direct control path).",
            "Changes the picture: 11.5ms of pure delay doesn't consume 90° of phase until ~22Hz, not ~6Hz — real headroom remains in the 10-20Hz band. The dominant constraint is more likely the plant's own ~38.5Hz resonance (visibly ringing in this data) than raw latency.",
        ],
        RESULTS + "fta_loop_delay_slide_summary.png",
        img_left=Inches(6.3),
    )

    add_full_image_slide(
        prs,
        "Does throttling help because it's more REGULAR, not just less frequent?",
        [
            "Open question: throttling to 200Hz recovered old Kp=1.75/Ki=200 behavior, but why is still unexplained — sample rate itself, or how UNIFORM the control-update timing is?",
            "Measured directly via dac_y as a firing-detector (apply_dac only runs inside a real control step), firmware-clock timestamped throughout: full rate CV (relative jitter) ≈ 71% (long tail to 20+ms) vs. throttled-200Hz CV ≈ 44% (tighter, shorter tail) — genuinely more regular, not just slower, though n=21 for the throttled case is a small sample and 44% CV is still real jitter.",
        ],
        RESULTS + "fta_ctrl_jitter_check_final.png",
    )

    add_full_image_slide(
        prs,
        "Tested a boxcar smoothing pre-filter — first result was a false alarm, real test says no",
        [
            "Full-rate + smoothing on Kp=1.75/Ki=200 looked like a clean win (161ms/3.5%/247ms) — until re-running full-rate WITHOUT smoothing on the exact same config also came back clean. The original instability had stopped reproducing for unrelated reasons — same mystery as the throttle recovery, not a smoothing effect (also makes sense mathematically: at full rate the averaging window is ~1 sample, a near no-op).",
            "Found a config that's still genuinely unstable right now (Kp=2.5/Ki=15) to test properly: smoothing does NOT rescue it at full rate (still 195% overshoot). Throttling alone DOES stabilize it (slow, ~2.2s). Adding smoothing on top of throttling changes nothing measurable — visually identical traces.",
            "Honest conclusion: rate reduction (throttling) is the real lever; boxcar averaging adds no measurable benefit on top of it. The deeper mystery — why the original instability stopped reproducing at all — remains open and is arguably the more important unresolved question now.",
        ],
        RESULTS + "fta_smoothing_vs_throttle_comparison.png",
    )

    # closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Where this leaves us")
    _add_bullets(
        slide,
        [
            "Best working point found: Kp=1.75 / Ki=200, throttled to 200Hz — settling ~15x faster and tracking gain ~2-3x higher (5-20Hz) than the un-throttled Ki=15 fallback.",
            "But even this doesn't solve the mission: sine validation shows S (disturbance sensitivity) still 66-80% across 5-20Hz — most of a real beacon wobble would still show up as position error. Fast step settling ≠ good disturbance rejection.",
            "Real closed-loop delay measured directly at ~11.5ms (not the ~41ms this project had assumed) — more phase-margin headroom than feared, so the resonance (~38.5Hz), not raw latency, is the more likely dominant constraint.",
            "Boxcar smoothing tested and ruled out — throttling's rate reduction is the real lever, averaging adds nothing on top of it.",
            "The single most important open question: the ORIGINAL Kp=1.75/Ki=200 full-rate instability that motivated this whole day's investigation has stopped reproducing, for reasons unconnected to any software change made today. This casts real doubt on how durable ANY of today's \"stable\" configurations actually are — worth understanding before trusting them long-term.",
            "Notch filter: implemented, hardware-validated, doesn't move the Ki/Kp boundary — a real negative result. D-term: ruled out at every setting tried, but only under a since-fixed throttle bug — worth one clean retest.",
        ],
        Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.9), size=13.5,
    )

    out = "docs/session_results_2026-08-18_pid_tuning.pptx"
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    main()
